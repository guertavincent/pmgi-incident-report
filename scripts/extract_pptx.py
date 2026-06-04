import json
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation

PPTX_PATH = Path(r"c:\Users\PMGO170326C\Documents\PMGI INCIDENT REPORT SYSTEM\pmgi-incident-report\PMGI-OPS-GL-01 - Houskeeping Training Guidelines rev.1.pptx")
OUTPUT_DIR = Path(r"c:\Users\PMGO170326C\Documents\PMGI INCIDENT REPORT SYSTEM\pmgi-incident-report\public\housekeeping-guidelines")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

prs = Presentation(str(PPTX_PATH))
slides_data = []


SPEAKER_AHASH = "0000010000111110111100101111001111110010001100100001010000000000"


def ahash_bits(img: Image.Image, size: int = 8) -> list[int]:
    img = img.convert("L").resize((size, size))
    data = list(img.getdata())
    avg = sum(data) / len(data)
    return [1 if p >= avg else 0 for p in data]


def is_speaker_icon(blob: bytes) -> bool:
    try:
        img = Image.open(BytesIO(blob))
    except Exception:
        return False

    bits = ahash_bits(img)
    base = [1 if c == "1" else 0 for c in SPEAKER_AHASH]
    dist = sum(1 for a, b in zip(bits, base) if a != b)
    return dist <= 6

for s_idx, slide in enumerate(prs.slides, start=1):
    slide_entry = {
        "index": s_idx,
        "texts": [],
        "images": [],
    }
    image_counter = 0

    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                slide_entry["texts"].append(text)

        if shape.shape_type == 13:  # PICTURE
            image_counter += 1
            image = shape.image
            if is_speaker_icon(image.blob):
                continue
            ext = image.ext or "bin"
            image_name = f"slide-{s_idx:02d}-img-{image_counter:02d}.{ext}"
            image_path = OUTPUT_DIR / image_name
            with open(image_path, "wb") as f:
                f.write(image.blob)
            slide_entry["images"].append(f"/housekeeping-guidelines/{image_name}")

    slides_data.append(slide_entry)

json_path = OUTPUT_DIR / "slides.json"
with open(json_path, "w", encoding="utf-8") as f:
    json.dump(slides_data, f, indent=2)

print(f"Extracted {len(slides_data)} slides to {OUTPUT_DIR}")
